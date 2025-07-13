#import tkinter as tk
#from tkinter import Tk, filedialog
import streamlit as st
import io
import os
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook
#import tkinter as tk
#from tkinter import filedialog
#import pandas as pd
from pptx import Presentation
import os
#import python-pptx as pptx
from pptx import Presentation
from pptx.dml.color import RGBColor
import streamlit as st
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
#import pptx
import pandas as pd
from pandas import read_excel
from openpyxl import load_workbook
from pptx.util import Inches
import openpyxl
from matplotlib.pyplot import savefig, subplots
import io
from PIL import Image
from pptx.enum.chart import XL_LABEL_POSITION
#import pandas as pd
from IPython.display import display
import ipywidgets as widgets
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.xmlchemy import OxmlElement
#from pptx.enum.table import WD_TABLE_STYLE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.dml import MSO_FILL
from datetime import datetime
import matplotlib.pyplot as plt
import zipfile
import math

def count_images_in_folder(folder_path):
    image_files = [f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    return len(image_files)

def count_images_in_zip(zip_buffer):
    with zipfile.ZipFile(zip_buffer, 'r') as zip_ref:
        image_files = [f for f in zip_ref.namelist() if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    return len(image_files)

def extract_images_from_zip(zip_buffer, extract_to_folder):
    with zipfile.ZipFile(zip_buffer, 'r') as zip_ref:
        zip_ref.extractall(extract_to_folder)
    return [f for f in os.listdir(extract_to_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

def read_excel_and_write_to_pptx(excel_path, pptx_path , image_folder_path):

    #image_zip_file_path = "path_to_uploaded_zip_file.zip"
    #image_folder_path = "extracted_images"

    # Extract images from the zip file
    #image_files = extract_images_from_zip(image_zip_file_path, image_folder_path)

    # Count the number of images in the extracted folder
    #num_images = count_images_in_zip(image_zip_file_path)

    image_files = [f for f in os.listdir(image_folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    # Count the number of images in the folder
    num_images = count_images_in_folder(image_folder_path)



    #df = read_excel(excel_path)
    ppt = Presentation(pptx_path)


    # Load the workbook using openpyxl
    excel_file = excel_path
    wb = load_workbook(excel_file, data_only=True)
    sheet = wb["CM"]
    sheet2 = wb["PM"] ### test PM tab in same code

    #df_cm = pd.read_excel(excel_path, sheet_name='CM', engine='openpyxl')
    #df_pm = pd.read_excel(excel_path, sheet_name='PM', engine='openpyxl')

    # Function to check if a row or column is hidden
    def is_hidden_row_or_column(sheet, row_idx=None, col_idx=None):
        if row_idx is not None:
            return sheet.row_dimensions[row_idx].hidden
        if col_idx is not None:
            col_letter = openpyxl.utils.get_column_letter(col_idx + 1)
            return sheet.column_dimensions[col_letter].hidden
        return False

    # Find the row containing 'Include in PPT' (header row)
    include_in_ppt_idx = None
    for row_idx, row in enumerate(sheet.iter_rows()):
        for cell in row:
            if cell.value and "include in ppt" in str(cell.value).lower():
                include_in_ppt_idx = row_idx
                break
        if include_in_ppt_idx is not None:
            break

    if include_in_ppt_idx is None:
        st.error(f"Could not find the header row for data in the 'CM' tab. Make sure we are following the agreed format.")
        st.stop()
        #raise ValueError("Could not find the header row for data in the 'CM' tab")

    # Access the header row values directly
    header_row = [cell.value for cell in sheet[include_in_ppt_idx + 1]]
    #print(f"header row: {header_row}")

    # Scan the header row for "Project Name"
    project_name_col = None
    include_in_ppt_col = None
    project_no_col = None
    site_area_col = None
    built_up_area_col = None
    cost_m2_col = None
    forecast_to_complete_col = None
    forecast_construction_spend_col = None
    payment_progress_col = None
    #include_in_ppt_col = None
    design_supervision_col = None
    civil_structure_finishes_col = None
    ffe_col = None
    services_col = None
    features_col = None
    external_works_col = None
    hs_landscaping_col = None
    hardscape_col = None
    softscape_col = None
    inf_utility_building_col = None
    inf_networks_col = None
    main_contractor_preliminaries_col = None
    #main_contractor_overhead_profit_col = None
    client_direct_procurement_col = None
    contigency_col = None
    variation_col = None
    vat_col = None

    current_project_cost_col = None
    project_no_pm_col = None
    project_status_col = None
    design_status_col = None
    construction_start_date_col = None
    target_completion_date_col = None
    forecast_completion_date_col = None
    overall_progress_col = None
    construction_progress_col = None
    remaining_col = None
    project_config_col = None
    space_programme_col = None
    daily_manpower_col = None
    daily_machinery_col = None
    remaining_items_col = None
    risk_assessment_col = None
    basis_project_col = None

    project_name_col_sheet2 = None

    project_name = ""
    include_in_ppt = ""
    project_no = ""
    site_area = ""
    built_up_area = ""
    cost_m2 = ""
    forecast_to_complete = ""
    forecast_construction_spend = ""
    payment_progress = ""
    #include_in_ppt_col = None
    design_supervision = ""
    civil_structure_finishes = ""
    ffe = ""
    services = ""
    features = ""
    external_works = ""
    hs_landscaping = ""
    hardscape = ""
    softscape = ""
    inf_utility_building = ""
    inf_networks = ""
    main_contractor_preliminaries = ""
    #main_contractor_overhead_profit_col = None
    client_direct_procurement = ""
    contigency = ""
    variation = ""
    vat = ""
    basis_project = ""

    project_name_PM = ""
    current_project_cost = ""
    project_status = ""
    design_status = ""
    construction_start_date = ""
    target_completion_date = ""
    forecast_completion_date = ""
    overall_progress = ""
    construction_progress = ""
    remaining = ""
    project_config = ""
    space_programme = ""
    daily_manpower = ""
    daily_machinery = ""
    remaining_items = ""
    risk_assessment = ""

    for idx, cell in enumerate(header_row): #capture data for relevant, non-blank rows
      if cell is not None:
        cell_value = str(cell).strip()
        #print(cell_value)
        if cell_value and not is_hidden_row_or_column(sheet, col_idx=idx):
            if "Project Name" in cell_value:
                project_name_col = idx
            elif "Basis of Project Cost" in cell_value:
                basis_project_col = idx
            elif "Include in PPT" in cell_value:
                include_in_ppt_col = idx
            elif "No." in cell_value:
                project_no_col = idx
            elif "Site Area" in cell_value:
                site_area_col = idx
            elif "Built up Area" in cell_value:
                built_up_area_col = idx
            elif "Cost / m2" in cell_value:
                cost_m2_col = idx
            elif "Forecast to Complete" in cell_value:
                forecast_to_complete_col = idx
            elif "Forecast Construction Spend" in cell_value:
                forecast_construction_spend_col = idx
            elif "Payment Progress" in cell_value:
                payment_progress_col = idx
            elif "Design & Supervision" in cell_value:
                design_supervision_col = idx
            elif "Civil, Structural" in cell_value:
                civil_structure_finishes_col = idx
            elif "FF&E" in cell_value:
                ffe_col = idx
            elif "Services" in cell_value:
                services_col = idx
            elif "Features" in cell_value:
                features_col = idx
            elif "External Works" in cell_value:
                external_works_col = idx
            elif "H&S Landscaping" in cell_value:
                hs_landscaping_col = idx
            elif "Hardscape" in cell_value:
                hardscape_col = idx
            elif "Softscape" in cell_value:
                softscape_col = idx
            elif "Inf. Utility Building" in cell_value:
                inf_utility_building_col = idx
            elif "Inf. Networks" in cell_value:
                inf_networks_col = idx
            elif "Main Contractor Preliminaries" in cell_value:
                main_contractor_preliminaries_col = idx
            elif "Main Contractor Overhead & Profit" in cell_value:
                main_contractor_overhead_profit_col = idx
            elif "Client Direct Procurement" in cell_value:
                client_direct_procurement_col = idx
            elif "Contingency" in cell_value:
                contigency_col = idx
            elif "Variation/Claims" in cell_value:
                variation_col = idx
            elif "VAT" in cell_value:
                vat_col = idx
            elif "Current Project Cost" in cell_value:
                current_project_cost_col = idx

    #print(f"Project Name Column: {project_name_col}")
    if project_name_col is None:
        st.error(f"Project Name is blank in one of the to-be-included in PPT rows in the Excel file. Please re-upload and try again")
        st.stop()

    def format_number(number):
        if number == "":
            return "0 M"
        elif number == "0":
            return "0 M"
        elif number is None:
            return ""
        else:
            number = float(number)
            if number >= 1_000_000_000:
                return f"{number / 1_000_000_000:.2f} B"
            else:
                return f"{number / 1_000_000:.2f} M"

    # Count number of total data rows in the Excel sheet
    #slides_count = 0
    #for row in sheet.iter_rows(min_row=include_in_ppt_idx + 2):
    #    if not is_hidden_row_or_column(sheet, row_idx=row[0].row):
    #        if any(cell.value for cell in row):
    #            slides_count += 1

    #print(f"Number of data (project) rows: {slides_count}")
    #if slides_count == 0:
    #    st.error(f"Either the Excel file's CM tab has ALL hidden/blank rows or the 'Include in PPT column' is not set to 'yes' for any record. Please re-upload and try again")
    #    st.stop()




    #################################################----------------PM tab
    # Find the row containing 'Category' (header row)
    category_index = None
    for row_idx, row in enumerate(sheet2.iter_rows()):
        for cell in row:
            if cell.value and "category" in str(cell.value).lower():
                category_index = row_idx
                break
        if category_index is not None:
            break

    if category_index is None:
        st.error(f"Could not find the header row for data in the 'PM' tab. It should begin with 'Category'")
        st.stop()
        #raise ValueError("Could not find the header row for data in the 'CM' tab")

    # Access the header row values directly
    header_row_pm = [cell.value for cell in sheet2[category_index + 1]]
    #print(f"header row PM: {header_row_pm}")

    #slides_count_PM = 0
    #for row in sheet2.iter_rows(min_row=category_index + 2):
    #    if not is_hidden_row_or_column(sheet2, row_idx=row[0].row):
    #        if any(cell.value for cell in row):
    #            slides_count_PM += 1

    #print(f"Number of data (project) rows in PM tab: {slides_count_PM}")
    #if slides_count_PM == 0: #or slides_count != slides_count_PM:
    #    st.error(f"The Excel file's PM tab is missing or hidden project rows. Please fix, re-upload the Excel file and try again")
    #    st.stop()

    for idx, cell in enumerate(header_row_pm): #capture data for relevant, non-blank rows
      if cell is not None:
        cell_value = str(cell).strip()
        #print(cell_value)
        if cell_value and not is_hidden_row_or_column(sheet2, col_idx=idx):
            #if "Current Project Cost" in cell_value:
            #    current_project_cost_col = idx
            if "Project Status" in cell_value:
                project_status_col = idx
            elif "No." in cell_value:
                project_no_pm_col = idx
            elif "Project Name" in cell_value:
                project_name_col_sheet2 = idx
            elif "Design Status" in cell_value:
                design_status_col = idx
            elif "Construction Start Date" in cell_value:
                construction_start_date_col = idx
            elif "Target Completion Date" in cell_value:
                target_completion_date_col = idx
            elif "Forecast Completion Date" in cell_value:
                forecast_completion_date_col = idx
            elif "Overall Progress" in cell_value:
                overall_progress_col = idx
            elif "Construction Progress" in cell_value:
                construction_progress_col = idx
            elif "Remaining Progress" in cell_value:
                remaining_col = idx
            elif "Project Configuration" in cell_value:
                project_config_col = idx
            elif "Space Programme" in cell_value:
                space_programme_col = idx
            elif "Daily Manpower" in cell_value:
                daily_manpower_col = idx
            elif "Daily Machinery" in cell_value:
                daily_machinery_col = idx
            elif "Remaining Items" in cell_value:
                remaining_items_col = idx
            elif "Risk Assessment" in cell_value:
                risk_assessment_col = idx


    #-----------------------------------------------------------------LOOP starts here
    # Read the PPT file
    ppt = Presentation(pptx_path)
    slide_layout = ppt.slides[1].slide_layout  # layout

    # Extract the rows of actual data (ignoring hidden rows)
    data_rows = []
    for row_idx, row in enumerate(sheet.iter_rows(min_row=include_in_ppt_idx + 2, max_row=1000, values_only=True)):
        if not is_hidden_row_or_column(sheet, row_idx=row_idx + include_in_ppt_idx + 2):
            if row[include_in_ppt_col] == "yes" or row[include_in_ppt_col] == "Yes":  # if that particular row is to be included in the PPT or not
                data_rows.append(row)


    ######################### -------------------- PM Data array
    #data_rows_PM = []
    #for row_idx, row in enumerate(sheet2.iter_rows(min_row=category_index + 2, max_row=sheet2.max_row, values_only=True)):
    #    if not is_hidden_row_or_column(sheet2, row_idx=row_idx + category_index + 2):
    #            data_rows_PM.append(row)

    # Extract 'Project Name' values from data_rows
    project_names = {row[project_name_col] for row in data_rows}
    project_nos = {row[project_no_col] for row in data_rows}
    
    # Build a dictionary from the PM tab using Project Name as key (normalized)
    pm_data_map = {}
    for row_idx, row in enumerate(sheet2.iter_rows(min_row=category_index + 2, max_row=1000, values_only=True)):
        if not is_hidden_row_or_column(sheet2, row_idx=row_idx + category_index + 2):
            key = str(row[project_name_col_sheet2]).strip().lower() if row[project_name_col_sheet2] else None
            if key:
                pm_data_map[key] = row
    # Match PM rows based on Project Name in the same order as CM tab
    data_rows_PM = []
    for cm_row in data_rows:
        cm_project_name = str(cm_row[project_name_col]).strip().lower()
        matched_pm_row = pm_data_map.get(cm_project_name)
    
        if matched_pm_row:
            print(f"Match found for Project Name: {cm_project_name}")
            data_rows_PM.append(matched_pm_row)
        else:
            print(f"No PM data found for Project Name: {cm_project_name}")
            st.warning(f"No PM data found for Project Name: {cm_project_name}")
            data_rows_PM.append([None] * len(pm_data_map[next(iter(pm_data_map))]))  # Safe fallback for missing PM row






    # Populate data_rows_PM with matching rows from sheet2
    #data_rows_PM = []
    #for row_idx, row in enumerate(sheet2.iter_rows(min_row=category_index + 2, max_row=1000, values_only=True)):
    #    if not is_hidden_row_or_column(sheet2, row_idx=row_idx + category_index + 2):
    #        #if row[project_name_col_sheet2] in project_names:
    #        if row[project_no_pm_col] in project_nos:
    #            data_rows_PM.append(row)
    #            print(f"matching row/record:1")



    #if len(data_rows) != len(data_rows_PM):
    #    st.error(f"The Excel file's PM tab is missing or hidden project rows. Please fix, re-upload the Excel file and try again")
    #    st.stop()

    ##########---------------- FOR LOOP dynamic ------------------------------------###########
    ##########-----all code below in this function needs to be dynamic and looped------#########

    # Insert the new slide immediately starting from slide 3 (position 4)
    insert_position = 2  # Start inserting from slide 3 (index starts at 0)
    for data_row, pm_row in zip(data_rows, data_rows_PM): # Only produce the no. of slides as data rows in the Excel sheet
        new_slide = ppt.slides.add_slide(slide_layout)
        xml_slides = ppt.slides._sldIdLst  # Access the low-level XML structure of slides
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])  # Remove the slide that was just added
        xml_slides.insert(insert_position, slides[-1])  # Insert the slide at the current position

#### - match image name with project name,
#### - find image dimensions in OG
#### - upload in each iteration, or do it separately in another loop

        # Extract the first row of actual data (ignoring hidden rows)
        for row_idx, row in enumerate(sheet.iter_rows(min_row=include_in_ppt_idx + 2, max_row=1000, values_only=True)):
            if not is_hidden_row_or_column(sheet, row_idx=row_idx + include_in_ppt_idx + 2):
                if row[include_in_ppt_col] == "yes": # if that particular row is to be included in the PPT or not
                    if project_name_col is not None:
                        project_name = data_row[project_name_col]
                    if basis_project_col is not None:
                        basis_project = data_row[basis_project_col]
                    if include_in_ppt_col is not None:
                        include_in_ppt = data_row[include_in_ppt_col]
                    if site_area_col is not None:
                        site_area = data_row[site_area_col]
                    if built_up_area_col is not None:
                        built_up_area = data_row[built_up_area_col]
                    if cost_m2_col is not None:
                        cost_m2 = data_row[cost_m2_col]
                    if forecast_to_complete_col is not None:
                        forecast_to_complete = data_row[forecast_to_complete_col]
                    if forecast_construction_spend_col is not None:
                        forecast_construction_spend = data_row[forecast_construction_spend_col]
                    if payment_progress_col is not None:
                        payment_progress = data_row[payment_progress_col]
                    if design_supervision_col is not None:
                        design_supervision = data_row[design_supervision_col]
                    if civil_structure_finishes_col is not None:
                        civil_structure_finishes = data_row[civil_structure_finishes_col]
                    if ffe_col is not None:
                        ffe = data_row[ffe_col]
                    if services_col is not None:
                        services = data_row[services_col]
                    if features_col is not None:
                        features = data_row[features_col]
                    if external_works_col is not None:
                        external_works = data_row[external_works_col]
                    if hs_landscaping_col is not None:
                        hs_landscaping = data_row[hs_landscaping_col]
                    if hardscape_col is not None:
                        hardscape = data_row[hardscape_col]
                    if softscape_col is not None:
                        softscape = data_row[softscape_col]
                    if inf_utility_building_col is not None:
                        inf_utility_building = data_row[inf_utility_building_col]
                    if inf_networks_col is not None:
                        inf_networks = data_row[inf_networks_col]
                    if main_contractor_preliminaries_col is not None:
                        main_contractor_preliminaries = data_row[main_contractor_preliminaries_col]
                    if client_direct_procurement_col is not None:
                        client_direct_procurement = data_row[client_direct_procurement_col]
                    if contigency_col is not None:
                        contigency = data_row[contigency_col]
                    if variation_col is not None:
                        variation = data_row[variation_col]
                    if vat_col is not None:
                        vat = data_row[vat_col]
                    if current_project_cost_col is not None:
                        current_project_cost = data_row[current_project_cost_col]
                #break
        #print(f"Project Name: {project_name} and Include in PPT:{include_in_ppt}")

        # Extract the first row of actual PM TAB DATA.. (ignoring hidden rows)
        for row_idx, row in enumerate(sheet2.iter_rows(min_row=category_index + 2, max_row= 1000, values_only=True)):
            if not is_hidden_row_or_column(sheet2, row_idx=row_idx + category_index + 2):
                #if current_project_cost_col is not None:
                #    current_project_cost = pm_row[current_project_cost_col]
                if project_status_col is not None:
                    project_status = pm_row[project_status_col]
                if design_status_col is not None:
                    design_status = pm_row[design_status_col]
                if construction_start_date_col is not None:
                    construction_start_date = pm_row[construction_start_date_col]
                if target_completion_date_col is not None:
                    target_completion_date = pm_row[target_completion_date_col]
                if forecast_completion_date_col is not None:
                    forecast_completion_date = pm_row[forecast_completion_date_col]
                if overall_progress_col is not None:
                    overall_progress = pm_row[overall_progress_col]
                if construction_progress_col is not None:
                    construction_progress = pm_row[construction_progress_col]
                if remaining_col is not None:
                    remaining = pm_row[remaining_col]
                    #print(f"remaining percentage:{remaining}")
                if project_config_col is not None:
                    project_config = pm_row[project_config_col]
                if space_programme_col is not None:
                    space_programme = pm_row[space_programme_col]
                if daily_manpower_col is not None:
                    daily_manpower = pm_row[daily_manpower_col]
                if daily_machinery_col is not None:
                    daily_machinery = pm_row[daily_machinery_col]
                if remaining_items_col is not None:
                    remaining_items = pm_row[remaining_items_col]
                if risk_assessment_col is not None:
                    risk_assessment = pm_row[risk_assessment_col]

        #print(f"project_status:{project_status}")

        # Process each image file
        for image_file in image_files:
            image_name, _ = os.path.splitext(image_file)
            if str(project_name).lower() == str(image_name).lower():
                image_path = os.path.join(image_folder_path, image_file)
                new_slide.shapes.add_picture(image_path, left=320893, top=4800600, width=4858788, height=4074726)


        for shape in new_slide.shapes:
        # PLACEHOLDER slide heading text (no. + Project Name) - this shape comes bydefault when inserting a new slide in such a layoutt
            if shape.left == 323167 and shape.top == 320054:
                text_frame = shape.text_frame
                text_frame.clear()  # Clear existing text (if any)
                # Add project_name to the placeholder shape
                p = text_frame.paragraphs[0]
                p.text = str(project_name)
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                p.font.size = Pt(24)  # Set font size (optional)
                p.alignment = PP_ALIGN.RIGHT  # Align text to the right
                p.font.name = 'Tajawal'

                # Insert a textbox with slide_count to the very left inside the placeholder shape
                #left = shape.left - Inches(1.0)  # Adjust position to the left of the placeholder shape
                left = 323166
                #top = shape.top
                top = 432136
                #width = Inches(1.0)
                width = 451534
                #height = shape.height
                height = 461665

                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                p1 = text_frame.paragraphs[0]
                #text_frame.clear()
                p1.text = str(insert_position-1)
                insert_position += 1  # Update the position for the next slide/iteration
                p1.font.size = Pt(24)  # Set font size (optional)
                p1.alignment = PP_ALIGN.LEFT  # Align text to the left
                p1.font.name = 'Tajawal'

        # Define properties for the first table (table1) - always static
        table1_position = (304800, 1125677)
        table1_height = 1166985
        table1_width = 12173635
        table1_rows = 2
        table1_columns = 10
        first_row_color = RGBColor(21, 42, 93)  # RGB(21, 42, 93)
        second_row_first_column_color = RGBColor(233, 245, 245)  # RGB(233, 245, 245)

        # Add the first table to the new slide
        table1 = new_slide.shapes.add_table(
            rows=table1_rows,
            cols=table1_columns,
            left=table1_position[0],
            top=table1_position[1],
            width=table1_width,
            height=table1_height
        )

        # Row heights adjustment
        row_heights = [int(table1_height * 0.45), int(table1_height * 0.45)]  # 15% less than original height
        for i, row in enumerate(table1.table.rows):
            row.height = row_heights[i]

        # Set headers for the first row (this is static and common for table1 across slides)
        headers = [
            "Project Name",
            "Project Status",
            "Design Status",
            "Construction Start Date",
            "Target Completion Date",
            "Forecast Completion Date",
            "Overall Progress",
            "Current Project Cost",
            "Forecast to Complete",
            "Cost / m2"
        ]

        first_row_cells = table1.table.rows[0].cells
        for cell in first_row_cells:
            cell.margin_left = 0
            cell.margin_right = 0
            cell.margin_top = 0
            cell.margin_bottom = 0
        second_row_second_columns = table1.table.rows[1]
        for i, cell in enumerate(second_row_second_columns.cells):
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(233, 245, 245)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        table1.table.cell(1, 0).text = str(project_name)
        table1.table.cell(1,0).text_frame.paragraphs[0].font.bold = True
        table1.table.cell(1,0).text_frame.paragraphs[0].font.name = 'Tajawal'
        table1.table.cell(1,0).text_frame.paragraphs[0].font.size = Pt(12)
        table1.table.cell(1, 1).text = str(project_status) if project_status is not None else ""
        table1.table.cell(1, 2).text = str(design_status) if design_status is not None else ""




        if target_completion_date is None:
            table1.table.cell(1, 4).text = ""
        elif isinstance(target_completion_date, datetime):
            table1.table.cell(1, 4).text = target_completion_date.strftime("%d %b %Y")
        else:
            table1.table.cell(1, 4).text = str(target_completion_date)

        if construction_start_date is None:
            table1.table.cell(1, 3).text = ""
        elif isinstance(construction_start_date, datetime):
            table1.table.cell(1, 3).text = construction_start_date.strftime("%d %b %Y")
        else:
            table1.table.cell(1, 3).text = str(construction_start_date)

        if forecast_completion_date is None:
            table1.table.cell(1, 5).text = ""
        elif isinstance(forecast_completion_date, datetime):
            table1.table.cell(1, 5).text = forecast_completion_date.strftime("%d %b %Y")
        else:
            table1.table.cell(1, 5).text = str(forecast_completion_date)


        if overall_progress is None:
            table1.table.cell(1, 6).text = ""
        elif isinstance(overall_progress, (int, float)):
            table1.table.cell(1, 6).text = f"{overall_progress * 100:.0f}%"
        else:
            table1.table.cell(1, 6).text = str(overall_progress)


        table1.table.cell(1,7).text = "SAR " + str(format_number(current_project_cost))
        table1.table.cell(1,8).text = "SAR " + str(format_number(forecast_to_complete))
        if cost_m2 is None:
            table1.table.cell(1, 9).text = ""
        elif isinstance(cost_m2, str):
            if cost_m2.startswith("#"):
                table1.table.cell(1, 9).text = " SAR / m²"
            else:
                table1.table.cell(1, 9).text = cost_m2 + " SAR / m²"
        elif isinstance(cost_m2, (int, float)):
            table1.table.cell(1, 9).text = str("{:,}".format(round(cost_m2))) + " SAR / m²"
        else:
            table1.table.cell(1, 9).text = str(cost_m2) + " SAR / m²"


        # Style the first row with the given color and white, bold text
        first_row = table1.table.rows[0]
        for i, cell in enumerate(first_row.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = first_row_color  # Set the background color for first row cells
            cell.text_frame.text = headers[i]  # Add header text to the cell
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
            cell.text_frame.paragraphs[0].font.bold = True  # Make text bold
            cell.text_frame.paragraphs[0].font.size = Pt(13)  # Optional: Adjust the font size if needed
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text
            cell.text_frame.paragraphs[0].font.name = 'Tajawal'
            cell.text_frame.margin_top = Pt(0)  # Set top margin to 0
            cell.text_frame.margin_bottom = Pt(0)  # Set bottom margin to 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center align the text

        second_row = table1.table.rows[1]
        for i, cell in enumerate(second_row.cells):
            cell.fill.solid()
            #cell.fill.fore_color.rgb = first_row_color  # Set the background color for first row cells
            #cell.text_frame.text = headers[i]  # Add header text to the cell
            #cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
            #cell.text_frame.paragraphs[0].font.bold = True  # Make text bold
            cell.text_frame.paragraphs[0].font.size = Pt(13)  # Optional: Adjust the font size if needed
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text
            cell.text_frame.paragraphs[0].font.name = 'Tajawal'
            cell.text_frame.margin_top = Pt(0)  # Set top margin to 0
            cell.text_frame.margin_bottom = Pt(0)  # Set bottom margin to 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center align the text

        # Set column widths
        column_widths = [1533857, 1071297, 1062459, 1202921, 1394530, 1372944, 867980, 1416027, 1164051, 1085155]
        for i, width in enumerate(column_widths):
            table1.table.columns[i].width = width

        for row in table1.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        # Define properties for the second table (table2)
        table2_position = (304800, 2271153)
        table2_height = 2287085
        table2_width = 4784621
        table2_first_row_color = RGBColor(21, 42, 93)  # RGB(21, 42, 93)
        table2_second_to_sixth_row_color = RGBColor(233, 245, 245)  # RGB(233, 245, 245)

        # Add the second table to the new slide
        table2 = new_slide.shapes.add_table(
            rows=5,
            cols=3,
            left=table2_position[0],
            top=table2_position[1],
            width=table2_width,
            height=table2_height
        )

        # Set first row merged with the text "Project Key Stats"
        first_row_table2 = table2.table.rows[0]
        for i, cell in enumerate(first_row_table2.cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(21, 42, 93)  # Set the background color for first row cells

        # Merge first row cells (if necessary)
        for i in range(2, 3):
            table2.table.cell(0, i).merge(table2.table.cell(0, 0))
        first_row_table2.cells[0].text_frame.text = "Project Key Stats"  # Add header text
        first_row_table2.cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        first_row_table2.cells[0].text_frame.paragraphs[0].font.bold = True  # Make text bold
        first_row_table2.cells[0].text_frame.paragraphs[0].font.size = Pt(13)  # Font size
        first_row_table2.cells[0].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
        first_row_table2.cells[0].text_frame.paragraphs[0].font.name = 'Tajawal'

        # Set the background color for the first two columns (rows 2 through 5)
        for row_idx in range(1, 5):
            for col_idx in range(3):  # For first two columns
                if col_idx == 2:
                    cell = table2.table.cell(row_idx, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                else:
                    cell = table2.table.cell(row_idx, col_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = table2_second_to_sixth_row_color

        # Set merged rows 2-4 with respective text
        merged_texts = ["Site Area", "Built Up Area", "Project Configuration"]
        for i, text in enumerate(merged_texts):
            cell1 = table2.table.cell(i + 1, 0)
            cell2 = table2.table.cell(i + 1, 1)
            cell1.merge(cell2)  # Merge cells in the first two columns
            cell1.text_frame.text = text  # Set the text
            cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
            cell1.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to Pt(12)
            cell1.text_frame.paragraphs[0].font.bold = True  # Bold text
            cell1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
            cell1.text_frame.paragraphs[0].font.name = 'Tajawal'

        # Merge row 5's columns 1 and 2 with text "Space Programme"
        cell5 = table2.table.cell(4, 0)
        cell6 = table2.table.cell(4, 1)
        cell5.merge(cell6)  # Merge cells in row 5, columns 1 and 2
        cell5.text_frame.text = "Space Programme"  # Add text
        cell5.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
        cell5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
        cell5.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to Pt(12)
        cell5.text_frame.paragraphs[0].font.bold = True  # Bold text
        cell5.text_frame.paragraphs[0].font.name = 'Tajawal'

        #table2.table.cell(1,1).text = str(

        # Adjust column widths for table2
        table2.table.columns[0].width = 999725
        table2.table.columns[1].width = 999725
        table2.table.columns[2].width = 2785175
        table2.table.rows[0].height = 320050
        table2.table.rows[1].height = 281850
        table2.table.rows[2].height = 281850
        table2.table.rows[3].height = 281850
        table2.table.rows[4].height = 1121475  # Adjusted height for merged row

        if site_area is None:
            table2.table.cell(1, 2).text = " m²"
        elif isinstance(site_area, (int, float)):
            table2.table.cell(1, 2).text = "{:,}".format(site_area) + ' m²'
        else:
            table2.table.cell(1, 2).text = str(site_area) + ' m²'
        #print(f"site area:{site_area}")
        table2.table.cell(1,2).text_frame.paragraphs[0].font.size = Pt(10)
        table2.table.cell(1,2).text_frame.paragraphs[0].font.name = 'Tajawal'
        table2.table.cell(1,2).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text


        # Check if built_up_area is a string or blank, and output it as is if it is
# Check if built_up_area is a string, blank, or None, and output it as is if it is
        if built_up_area is None:
            table2.table.cell(2, 2).text = " m²"
        elif isinstance(built_up_area, (int, float)):
            table2.table.cell(2, 2).text = "{:,}".format(built_up_area) + ' m²'
        else:
            table2.table.cell(2, 2).text = str(built_up_area) + ' m²'
        table2.table.cell(2,2).text_frame.paragraphs[0].font.size = Pt(10)
        table2.table.cell(2,2).text_frame.paragraphs[0].font.name = 'Tajawal'
        table2.table.cell(2,2).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text


        project_config = str(project_config if project_config is not None else "")
        table2.table.cell(3, 2).text = project_config
        table2.table.cell(3,2).text_frame.paragraphs[0].font.size = Pt(10)
        table2.table.cell(3,2).text_frame.paragraphs[0].font.name = 'Tajawal'
        table2.table.cell(3,2).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text


        if space_programme:
            lines = space_programme.split('\n')
        else:
            lines = []

        # Get the text frame of the cell
        text_frame = table2.table.cell(4, 2).text_frame

        # Clear any existing paragraphs
        text_frame.clear()

        if lines:
            # Add each line as a bullet point
            for line in lines:
                p = text_frame.add_paragraph()
                p.text = f"- {line}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text color
                p.level = 0  # Bullet point level
                p.alignment = PP_ALIGN.LEFT

            # Optionally, set the font name for all paragraphs
            for p in text_frame.paragraphs:
                p.font.name = "Tajawal"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT
        else:
            text_frame.text = ""

        table2.table.cell(4,2).text_frame.paragraphs[0].font.size = Pt(10)
        table2.table.cell(4,2).text_frame.paragraphs[0].font.name = 'Tajawal'
        table2.table.cell(4,2).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text


        for row in table2.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        # Add the third table (table3) with 10 rows and 2 columns ############--------------------------------this needs to be DYNAMIC-------------------------####
        table3_position = (5344484, 3087191)
        table3_height = 2324670
        table3_width = 2384824
        table3_rows = 10
        table3_columns = 2

        table3 = new_slide.shapes.add_table(
            rows=table3_rows,
            cols=table3_columns,
            left=table3_position[0],
            top=table3_position[1],
            width=table3_width,
            height=table3_height
        ).table

        # Set background color and black non-bold text for Table 3
        for row in range(table3_rows):
            for col in range(table3_columns):
                cell = table3.cell(row, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set cell background color to white
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                cell.text_frame.paragraphs[0].font.bold = False  # Non-bold text
                cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size to 12
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text

                ## Set borders for each cell
                #for border in ['left', 'right', 'top', 'bottom']:
                    #line = cell.border(border)
                    #line.color.rgb = RGBColor(59, 134, 134)  # Set border color
                    #line.width = Pt(1)  # Set border width


        # Textboxes (Textbox1 and Textbox2)
        # First Textbox
        textbox1 = new_slide.shapes.add_textbox(
            left=7900000,
            top=2271425,
            width=4590000,
            height=348040
        )
        textbox1.fill.solid()
        textbox1.fill.fore_color.rgb = RGBColor(21, 42, 93)  # Set background color
        text_frame1 = textbox1.text_frame
        text_frame1.text = "Current Construction & Payment Progress"
        text_frame1.paragraphs[0].font.size = Pt(13)
        text_frame1.paragraphs[0].font.bold = True
        text_frame1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        text_frame1.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame1.paragraphs[0].font.name = 'Tajawal'

        # Second Textbox
        textbox2 = new_slide.shapes.add_textbox(
            left=5337406,
            top=2271425,
            width=2375000,
            height=348040
        )
        textbox2.fill.solid()
        textbox2.fill.fore_color.rgb = RGBColor(21, 42, 93)  # Set background color
        text_frame2 = textbox2.text_frame
        text_frame2.text = "Project Cost & Cost Analysis"
        text_frame2.paragraphs[0].font.size = Pt(13)
        text_frame2.paragraphs[0].font.bold = True
        text_frame2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        text_frame2.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame2.paragraphs[0].font.name = 'Tajawal'

       # Shape 11: Table
        shape_11 = new_slide.shapes.add_table(2, 1, Inches(7626125 / 914400), Inches(6543793 / 914400), Inches(4852309 / 914400), Inches(1238901 / 914400)).table

        # Set row heights (30% for row 1 and 70% for row 2)
        total_height = int(Inches(1179842 / 914400))
        shape_11.rows[0].height = int(total_height * 0.26)
        shape_11.rows[1].height = int(total_height * 0.74)
        shape_11.cell(0, 0).text = "Remaining Items"

        if remaining_items:
            lines = remaining_items.split('\n')
        else:
            lines = []

        # Get the text frame of the cell
        text_frame = shape_11.cell(1, 0).text_frame

        # Clear any existing paragraphs
        text_frame.clear()

        if lines:
            # Add each line as a bullet point
            for line in lines:
                p = text_frame.add_paragraph()
                p.text = f"- {line}"
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text color
                p.level = 0  # Bullet point level
                p.alignment = PP_ALIGN.LEFT
                p.bullet = True  # Enable bullet points
                p.bullet_char = '\u2022'  # Set bullet character to a solid circle

            # Optionally, set the font name for all paragraphs
            for p in text_frame.paragraphs:
                p.font.name = "Tajawal"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.LEFT
        else:
            text_frame.text = ""

        # Set the first row background color to RGBColor(21, 42, 93) with white bold text
        shape_11.cell(0, 0).fill.solid()
        shape_11.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
        cell_11_text_frame = shape_11.cell(0, 0).text_frame
        cell_11_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        cell_11_text_frame.paragraphs[0].font.bold = True
        cell_11_text_frame.paragraphs[0].font.size = Pt(13)
        cell_11_text_frame.paragraphs[0].font.name = "Tajawal"
        shape_11.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text
        shape_11.cell(0, 0).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Set the second row background color to RGBColor(233, 245, 245)
        shape_11.cell(1, 0).fill.solid()
        shape_11.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)

        # Adjust cell margins and text color
        for row in shape_11.rows:
            for cell in row.cells:
                cell.margin_left = Inches(0.05)
                cell.margin_right = 0
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = 0
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        shape_11.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        shape_11.cell(0,0).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Shape 12: Table
        shape_12 = new_slide.shapes.add_table(3, 2, Inches(5342198 / 914400), Inches(6543794 / 914400), Inches(2134930 / 914400), Inches(1217006 / 914400)).table
        shape_12.cell(0, 0).text = "Estimated Total Required Avg. Daily Resources"
        # Merge the first row
        cell_12_0_0 = shape_12.cell(0, 0)
        cell_12_0_1 = shape_12.cell(0, 1)
        cell_12_0_0.merge(cell_12_0_1)
        # Set the first row background color to RGBColor(21, 42, 93) with white bold text
        cell_12_0_0.fill.solid()
        cell_12_0_0.fill.fore_color.rgb = RGBColor(21, 42, 93)
        cell_12_text_frame = cell_12_0_0.text_frame
        cell_12_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        cell_12_text_frame.paragraphs[0].font.bold = True
        cell_12_text_frame.paragraphs[0].font.size = Pt(13)
        cell_12_text_frame.paragraphs[0].font.name = "Tajawal"
        cell_12_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Set the second row background color to RGBColor(233, 245, 245)
        shape_12.cell(1, 0).fill.solid()
        shape_12.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
        shape_12.cell(1, 1).fill.solid()
        shape_12.cell(1, 1).fill.fore_color.rgb = RGBColor(233, 245, 245)
        # Insert the values "Daily Manpower" and "Daily Machinery" in the second row with black bold point 10.5 text
        shape_12.cell(1, 0).text = "Daily Manpower"
        shape_12.cell(1, 1).text = "Daily Machinery"

        # Check if daily_manpower and daily_machinery are None, and make them zero if they are
        daily_manpower = str(daily_manpower if daily_manpower is not None else 0)
        daily_machinery = str(daily_machinery if daily_machinery is not None else 0)

        #shape_12.cell(1, 0).text = daily_manpower
        #shape_12.cell(1, 1).text = daily_machinery

        shape_12.cell(2, 0).text = str(daily_manpower)
        shape_12.cell(2, 1).text = str(daily_machinery)

        shape_12.cell(2, 0).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shape_12.cell(2, 1).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_12_text_frame_1 = shape_12.cell(1, 0).text_frame
        cell_12_text_frame_1.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        cell_12_text_frame_1.paragraphs[0].font.bold = True
        cell_12_text_frame_1.paragraphs[0].font.size = Pt(12)
        cell_12_text_frame_2 = shape_12.cell(1, 1).text_frame
        cell_12_text_frame_2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        cell_12_text_frame_2.paragraphs[0].font.bold = True
        cell_12_text_frame_2.paragraphs[0].font.size = Pt(12)
        cell_12_text_frame_1.paragraphs[0].font.name = "Tajawal"
        cell_12_text_frame_2.paragraphs[0].font.name = "Tajawal"
        cell_12_text_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
        cell_12_text_frame_2.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
        # Set the last row background color to white
        shape_12.cell(2, 0).fill.solid()
        shape_12.cell(2, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
        shape_12.cell(2, 1).fill.solid()
        shape_12.cell(2, 1).fill.fore_color.rgb = RGBColor(233, 245, 245)
        # Adjust row heights (30% for row 1 and 70% for row 2)
        total_height = Inches(1173428 / 914400)
        cell_12_text_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text

        for row in shape_12.rows:
            for cell in row.cells:
                cell.margin_left = 0
                cell.margin_right = 0
                cell.margin_top = 0
                cell.margin_bottom = 0
                cell.text_frame.paragraphs[0].font.size = Pt(13)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame = cell.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Set vertical alignment to middle

                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(13)
                    paragraph.alignment = PP_ALIGN.CENTER

                    # Set text direction to left-to-right (LTR)
                    paragraph._pPr.set('rtl', '0')

                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        shape_12.cell(2, 0).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shape_12.cell(2, 1).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        shape_12.cell(2, 0).margin_top = Inches(0.09)
        shape_12.cell(2, 1).margin_top = Inches(0.09)

        for paragraph in shape_12.cell(2, 0).text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph._pPr.set('rtl', '0')

        for paragraph in shape_12.cell(2, 1).text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph._pPr.set('rtl', '0')

        # Shape 13: Table
        shape_13 = new_slide.shapes.add_table(1, 2, Inches(5333250 / 914400), Inches(2673352 / 914400), Inches(2362950 / 914400), Inches(365760 / 914400)).table
        shape_13.cell(0, 0).text = "Current Project Cost"
        shape_13.cell(0, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        shape_13.cell(0, 0).fill.solid()
        shape_13.cell(0, 0).fill.fore_color.rgb = RGBColor(59, 134, 134)
        # Make cell (0,0)'s text bold and set font size to 12
        cell_00_text_frame = shape_13.cell(0, 0).text_frame
        shape_13.cell(0,1).text = str(format_number(current_project_cost))
        cell_00_text_frame.paragraphs[0].font.bold = True
        cell_00_text_frame.paragraphs[0].font.size = Pt(11)
        cell_00_text_frame.paragraphs[0].font.name = "Tajawal"
        shape_13.cell(0,1).text_frame.paragraphs[0].font.size = Pt(11)
        shape_13.cell(0,1).text_frame.paragraphs[0].font.name = "Tajawal"
        shape_13.cell(0,1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text color
        shape_13.cell(0,1).text_frame.paragraphs[0].font.bold = False
        # Set the second column's color
        shape_13.cell(0, 1).fill.solid()
        shape_13.cell(0, 1).fill.fore_color.rgb = RGBColor(211,236,236)
        # Set the column widths to 65% and 35% for the first two columns respectively
        total_width = int(Inches(2362950 / 914400))
        first_col_width = int(total_width * 0.67)
        second_col_width = int(total_width * 0.33)
        shape_13.columns[0].width = first_col_width
        shape_13.columns[1].width = second_col_width
        shape_13.rows[0].height =int(shape_13.rows[0].height*0.9)

        #shape_13.width = table3.width
        shape_13.width = Inches(2384824 / 914400)

        for row in shape_13.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        # Shape 15: Table
        shape_15 = new_slide.shapes.add_table(2, 1, Inches(5342195 / 914400), Inches(5567230 / 914400), Inches(7136239 / 914400), Inches(1034316 / 914400)).table
        shape_15.cell(0, 0).text = "Basis of Project Cost"
        # Set the first row background color to RGBColor(21, 42, 93) with white bold text
        shape_15.cell(0, 0).fill.solid()
        shape_15.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
        cell_15_text_frame = shape_15.cell(0, 0).text_frame
        cell_15_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        cell_15_text_frame.paragraphs[0].font.bold = True
        cell_15_text_frame.paragraphs[0].font.size = Pt(13)
        cell_15_text_frame.paragraphs[0].font.name = "Tajawal"
        # Set the second row background color to white
        shape_15.cell(1, 0).fill.solid()
        shape_15.cell(1, 0).fill.fore_color.rgb = RGBColor(255, 255, 255)
        for row in shape_15.rows:
            for cell in row.cells:
                #cell.margin_left = 0
                #cell.margin_right = 0
                #cell.margin_top = 0
                #cell.margin_bottom = 0
                cell.text_frame.paragraphs[0].font.size = Pt(13)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                #cell.text_frame.paragraphs[0].space_after = Pt(0)  # Minimal line spacing

                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        shape_15.rows[0].height = int(shape_15.rows[0].height * 0.6)

        if basis_project:
            lines = basis_project.split('\n')
        else:
            lines = []

        # Get the text frame of the cell
        text_frame = shape_15.cell(1, 0).text_frame

        # Clear any existing paragraphs
        text_frame.clear()

        if lines:
            # Add each line as a bullet point
            for line in lines:
                p = text_frame.add_paragraph()
                p.text = f"- {line}"
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0, 0, 0)  # Black text color
                p.level = 0  # Bullet point level
                p.alignment = PP_ALIGN.LEFT
                p.bullet = True  # Enable bullet points
                p.bullet_char = '\u2022'  # Set bullet character to a solid circle

            # Optionally, set the font name for all paragraphs
            for p in text_frame.paragraphs:
                p.font.name = "Tajawal"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.LEFT
        else:
            text_frame.text = ""
        
        values = []
        values2 = []

        if design_supervision != "" and design_supervision != 0:
            formatted_value = format_number(design_supervision)
            if formatted_value:
                values.append("Design & Supervision")
                values2.append(formatted_value)
        if civil_structure_finishes != "" and civil_structure_finishes != 0:
            formatted_value = format_number(civil_structure_finishes)
            if formatted_value:
                values.append("Civil, Structural")
                values2.append(formatted_value)
        if ffe != "" and ffe != 0:
            formatted_value = format_number(ffe)
            if formatted_value:
                values.append("FF&E")
                values2.append(formatted_value)
        if services != "" and services != 0:
            formatted_value = format_number(services)
            if formatted_value:
                values.append("Services")
                values2.append(formatted_value)
        if features != "" and features != 0:
            formatted_value = format_number(features)
            if formatted_value:
                values.append("Features")
                values2.append(formatted_value)
        if external_works != "" and external_works != 0:
            formatted_value = format_number(external_works)
            if formatted_value:
                values.append("External Works")
                values2.append(formatted_value)
        if hs_landscaping != "" and hs_landscaping != 0:
            formatted_value = format_number(hs_landscaping)
            if formatted_value:
                values.append("H&S Landscaping")
                values2.append(formatted_value)
        if hardscape != "" and hardscape != 0:
            formatted_value = format_number(hardscape)
            if formatted_value:
                values.append("Hardscape")
                values2.append(formatted_value)
        if softscape != "" and softscape != 0:
            formatted_value = format_number(softscape)
            if formatted_value:
                values.append("Softscape")
                values2.append(formatted_value)
        if inf_utility_building != "" and inf_utility_building != 0:
            formatted_value = format_number(inf_utility_building)
            if formatted_value:
                values.append("Inf. Utility Building")
                values2.append(formatted_value)
        if inf_networks != "" and inf_networks != 0:
            formatted_value = format_number(inf_networks)
            if formatted_value:
                values.append("Inf. Networks")
                values2.append(formatted_value)
        if main_contractor_preliminaries != "" and main_contractor_preliminaries != 0:
            formatted_value = format_number(main_contractor_preliminaries)
            if formatted_value:
                values.append("Preliminaries")
                values2.append(formatted_value)
        if client_direct_procurement != "" and client_direct_procurement != 0:
            formatted_value = format_number(client_direct_procurement)
            if formatted_value:
                values.append("Client Direct Procurement")
                values2.append(formatted_value)
        if contigency != "" and contigency != 0:
            formatted_value = format_number(contigency)
            if formatted_value:
                values.append("Contingency")
                values2.append(formatted_value)
        if variation != "" and variation != 0:
            formatted_value = format_number(variation)
            if formatted_value:
                values.append("Variation/Claims")
                values2.append(formatted_value)
        if vat != "" and vat != 0:
            formatted_value = format_number(vat)
            if formatted_value:
                values.append("VAT")
                values2.append(formatted_value)

        if not values:
            values.extend([
                "Design & Supervision", "Civil, Structural", "FF&E", "Services", "Features",
                "External Works", "H&S Landscaping", "Hardscape", "Softscape",
                "Inf. Utility Building", "Inf. Networks", "Preliminaries",
                "Client Direct Procurement", "Contingency", "Variation/Claims", "VAT"
            ])
            values2.extend([""] * 16)

        #if current_project_cost != "" and current_project_cost != 0:
        #    values.append("Current Project Cost")
        forecast_construction_spend = str(format_number(forecast_construction_spend))
        values.append(f"Forecast Construction Spend\n{forecast_construction_spend}")

        num_rows = len(values)
        num_cols = 2
        shape_20_width = int(Inches(2351715 / 914400))
        shape_20_height = int(Inches(2322108 / 914400))

        # Create shape_20 with the specified width and height
        shape_20 = new_slide.shapes.add_table(num_rows, num_cols, int(Inches(5344484 / 914400)), int(Inches(3087191 / 914400) + Inches(0.03)), shape_20_width, shape_20_height).table

        graphic_frame = shape_20._graphic_frame

        # Access the table's XML element
        tbl = graphic_frame.element.graphic.graphicData.tbl

        # Set the style ID
        style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
        tbl[0][-1].text = style_id
        #shape_20 = shape_20.table

        # Set the table design to "no style, no grid"
        # shape_20._graphic_frame.element.get_or_add_tblStyleLst().clear()

        # Set the first column width to 80% of the total width
        total_width = shape_20_width
        first_col_width = int(total_width * 0.7)
        second_col_width = int(total_width * 0.3)
        shape_20.columns[0].width = first_col_width
        shape_20.columns[1].width = second_col_width

        for row in shape_20.rows:
            for cell in row.cells:
                cell.margin_left = Inches(0.05)
                cell.margin_right = 0
                cell.margin_top = 0
                cell.margin_bottom = 0
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(211, 236, 236)  # Match shape13
                # cell = _set_cell_border(cell, border_color="FFFFFF", border_width='12700')
                # cell = apply_border(cell, edges=["left", "right", "top", "bottom"])
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                cell.text_frame.paragraphs[0].space_after = Pt(0)  # Minimal line spacing


        # Set the values in the first column and set text size to Pt(11) and not bold
        for i, value in enumerate(values):
            shape_20.cell(i, 0).text = value
            cell_text_frame = shape_20.cell(i, 0).text_frame
            cell_text_frame.paragraphs[0].font.size = Pt(11)
            cell_text_frame.paragraphs[0].font.bold = False
            cell_text_frame.paragraphs[0].font.name = "Tajawal"
            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        for i, value in enumerate(values2):
            shape_20.cell(i, 1).text = str(value)
            cell_text_frame = shape_20.cell(i, 1).text_frame
            cell_text_frame.paragraphs[0].font.size = Pt(11)
            cell_text_frame.paragraphs[0].font.bold = False
            cell_text_frame.paragraphs[0].font.name = "Tajawal"
            cell_text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Ensure the first row is not styled as a heading/header and has normal text like the other cells
        for i in range(num_cols):
            cell_text_frame = shape_20.cell(0, i).text_frame
            cell_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set text color to black
            cell_text_frame.paragraphs[0].font.bold = False
            cell_text_frame.paragraphs[0].font.size = Pt(11)
            cell_text_frame.paragraphs[0].font.name = "Tajawal"

        # Merge the cells in the last row
        shape_20.cell(num_rows - 1, 0).merge(shape_20.cell(num_rows - 1, 1))

        for row in shape_20.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                #cell.text_frame.paragraphs[0].space_after = Pt(0)  # Minimal line spacing
                cell_text_frame.paragraphs[0].font.name = "Tajawal"


        # Set the background color to RGB(29, 88, 137) and text to white bold for the last row
        last_row_cell = shape_20.cell(num_rows - 1, 0)
        last_row_cell.fill.solid()
        last_row_cell.fill.fore_color.rgb = RGBColor(29, 88, 137)
        last_row_text_frame = last_row_cell.text_frame

        # Clear any existing paragraphs
        last_row_text_frame.clear()

        # Set the background color to RGB(29, 88, 137) and text to white bold for the last row
        last_row_cell = shape_20.cell(num_rows - 1, 0)
        last_row_cell.fill.solid()
        last_row_cell.fill.fore_color.rgb = RGBColor(29, 88, 137)
        last_row_text_frame = last_row_cell.text_frame

        # Clear any existing paragraphs
        last_row_text_frame.clear()


        last_row_text_frame.text = "Forecast Construction Spend     " + "SAR " + str(forecast_construction_spend)

        last_row_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        last_row_text_frame.paragraphs[0].font.bold = True
        last_row_text_frame.paragraphs[0].font.name = "Tajawal"
        last_row_text_frame.paragraphs[0].font.size = Pt(11)
        last_row_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Ensure no extra line breaks are added
        last_row_text_frame.word_wrap = True

        # Adjust row heights to ensure the table fits within the specified height
        total_height = shape_20_height
        current_row_height = int(total_height / num_rows)

        # Set the height for each row
        for i in range(num_rows):
            shape_20.rows[i].height = current_row_height

        # Adjust the last row height to be 30% larger
        #new_last_row_height = int(current_row_height * 1.3)
        #height_difference = new_last_row_height - current_row_height
        #new_other_row_height = int((total_height - new_last_row_height) / (num_rows - 1))

        #for i in range(num_rows - 1):
        #    shape_20.rows[i].height = new_other_row_height

        #shape_20.rows[num_rows - 1].height = new_last_row_height

        for row in shape_20.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')

        # Shape 21: Table
        # Shape 21: Table
        shape_21 = new_slide.shapes.add_table(2, 1, Inches(5342195 / 914400), Inches(7841010 / 914400), Inches(7136239 / 914400), Inches(1034316 / 914400)).table
        shape_21.cell(0, 0).text = "Risk Assessment"

        # Set the first row background color to RGBColor(21, 42, 93) with white bold text
        shape_21.cell(0, 0).fill.solid()
        shape_21.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
        total_height = Inches(1015802.4 / 914400)
        cell_21_text_frame = shape_21.cell(0, 0).text_frame
        cell_21_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        cell_21_text_frame.paragraphs[0].font.bold = True
        cell_21_text_frame.paragraphs[0].font.size = Pt(13)
        cell_21_text_frame.paragraphs[0].font.name = "Tajawal"
        cell_21_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text

        # Set the second row background color to RGBColor(233, 245, 245)
        shape_21.cell(1, 0).fill.solid()
        shape_21.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)

        # Set row heights explicitly to avoid resizing
        shape_21.rows[0].height = int(total_height * 0.3)
        shape_21.rows[1].height = int(total_height * 0.7)

        #if risk_assessment:
        #    lines = risk_assessment.split('\n')
        #else:
        #    lines = []

        # Get the text frame of the cell
        #text_frame = shape_21.cell(1, 0).text_frame

        # Clear any existing paragraphs
        #text_frame.clear()

        #if lines:
            # Add each line as a bullet point without introducing extra line breaks
        #    for line in lines:
        #        p = text_frame.add_paragraph()
        #        #p.text = line.strip()  # Strip any leading/trailing whitespace
        #        p.text = f"- {line}"
        #        p.font.size = Pt(13)
        #        p.level = 0  # Bullet point level
        #        p.alignment = PP_ALIGN.LEFT

            # Optionally, set the font name for all paragraphs
        #    for p in text_frame.paragraphs:
        #        p.font.name = "Tajawal"
        #        p.font.color.rgb = RGBColor(0, 0, 0)
        #        p.font.size = Pt(13)
        #        p.alignment = PP_ALIGN.LEFT
        #else:
        #    text_frame.text = ""

        #text_frame = shape_21.cell(1, 0).text_frame
        #text_frame.clear()
        #lines = str(risk_assessment or "").strip().split('\n')
        
        #if lines:
        #    for idx, line in enumerate(lines):
        #        cleaned_line = line.strip()
        #        if idx == 0:
        #            text_frame.text = f"- {cleaned_line}"
        #        else:
        #            para = text_frame.add_paragraph()
        #            para.text = f"- {cleaned_line}"
        #        p = text_frame.paragraphs[idx]
        #        p.font.size = Pt(13)
        #        p.font.name = "Tajawal"
        #        p.font.color.rgb = RGBColor(0, 0, 0)
        #        p.alignment = PP_ALIGN.LEFT
        #else:
        #    text_frame.text = ""


        
        #text_frame = shape_21.cell(1, 0).text_frame
        #text_frame.clear()
        
        # Ensure the text frame is completely cleared
        #while text_frame.paragraphs:
        #    text_frame._element.clear_content()
        
        #lines = str(risk_assessment or "").strip().split('\n')
        
        #if lines:
        #    for idx, line in enumerate(lines):
        #        cleaned_line = line.strip()
        #        if cleaned_line:  # Ensure the line is not empty
        #            if idx == 0:
        #                text_frame.text = f"- {cleaned_line}"
        #            else:
        #                para = text_frame.add_paragraph()
        #                para.text = f"- {cleaned_line}"
        #            p = text_frame.paragraphs[idx]
        #            p.font.size = Pt(13)
        #            p.font.name = "Tajawal"
        #            p.font.color.rgb = RGBColor(0, 0, 0)
        #            p.alignment = PP_ALIGN.LEFT
        #else:
        #    text_frame.text = ""

        # Set the text color to white for the first row
        shape_21.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        for row in shape_21.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    #p._pPr.set('algn', 'l')
                    p._pPr.set('rtl', '0')
        ####################-====================================================================================DONUT 1

       # Ensure variables are not None, blank, or NaN
        #print(f"Construction Progress value:{construction_progress}")

        import matplotlib.pyplot as plt
        import io
        import math


        # Ensure variables are not None, blank, NaN, or Excel errors
        if construction_progress is None or (isinstance(construction_progress, float) and math.isnan(construction_progress)) or (isinstance(construction_progress, str) and construction_progress.startswith("#")):
            construction_progress = 0

        if remaining is None or (isinstance(remaining, float) and math.isnan(remaining)) or (isinstance(remaining, str) and remaining.startswith("#")):
            remaining = 0

        if payment_progress is None or (isinstance(payment_progress, float) and math.isnan(payment_progress)) or (isinstance(payment_progress, str) and payment_progress.startswith("#")):
            payment_progress = 0

        # If construction_progress is zero, make remaining 100
        if construction_progress == 0:
            remaining = 100

        # Data for the first donut chart
        sizes = [int(round(construction_progress * 100)), int(round(remaining * 100))]
        labels = ['Progress', 'Remaining']
        colors = ['#0aa57f', '#1d5889']  # RGB: (10, 165, 127) and (29, 88, 137)

        # Create a figure and axis with equal aspect ratio to avoid squeezing
        fig, ax = plt.subplots(figsize=(8, 8))


        sizes = [max(0, size) for size in sizes]

        # Create the donut chart with thicker white borders for slices
        wedges, texts, autotexts = ax.pie(sizes, colors=colors, autopct=lambda p: '{:.0f}%'.format(round(p)) if p > 7 else '',
                                          startangle=90, wedgeprops=dict(width=0.3, edgecolor='white', linewidth=3),
                                          pctdistance=0.85, textprops=dict(color='white', fontsize=28))

        # Add legends at the bottom of the chart in one line without borders or shadows and make them smaller horizontally
        ax.legend(wedges, labels, loc="upper center", bbox_to_anchor=(0.5, -0.1), frameon=False, ncol=2, fontsize=25,
                  handlelength=0.8, handleheight=0.8)

        # Equal aspect ratio ensures that pie is drawn as a circle
        ax.axis('equal')

        # Insert dynamic data into the center without borders
        ax.text(0, 0, 'Construction\n Progress:\n', ha='center', va='center', fontsize=30, fontname='Tajawal')
        ax.text(0, -0.2, str(int(round(construction_progress * 100))) + '%', ha='center', va='center', fontsize=40, fontname='Tajawal')

        # Save the plot to a BytesIO object
        buf = io.BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight')
        buf.seek(0)

        # Define the position and size of the image for the first chart
        left = 7884225
        top = 2964434
        width = 2245675
        height = 2551024  # Adjusted height to maintain aspect ratio

        # Add the image to the slide from BytesIO object
        new_slide.shapes.add_picture(buf, left, top, width, height)

        # If payment_progress is zero, make cost_to_complete 100
        if payment_progress == 0:
            cost_to_complete = 100
        else:
            cost_to_complete = int(round(100 - (payment_progress * 100)))

        # Data for the second donut chart
        sizes = [int(round(payment_progress * 100)), cost_to_complete]
        labels = ['Paid to Date', 'Cost to Complete']
        colors = ['#0aa57f', '#1d5889']  # RGB: (10, 165, 127) and (29, 88, 137)

        # Create a figure and axis with equal aspect ratio to avoid squeezing
        fig2, ax2 = plt.subplots(figsize=(8, 8))


        sizes = [max(0, size) for size in sizes]


        # Create the donut chart with thicker white borders for slices
        wedges2, texts2, autotexts2 = ax2.pie(sizes, colors=colors, autopct=lambda p: '{:.0f}%'.format(round(p)) if p > 7 else '',
                                              startangle=90, wedgeprops=dict(width=0.3, edgecolor='white', linewidth=3),
                                              pctdistance=0.85, textprops=dict(color='white', fontsize=28))

        # Add legends at the bottom of the chart in one line without borders or shadows and make them smaller horizontally
        ax2.legend(wedges2, labels, loc="upper center", bbox_to_anchor=(0.5, -0.1), frameon=False, ncol=2, fontsize=25,
                   handlelength=0.8, handleheight=0.8)

        # Equal aspect ratio ensures that pie is drawn as a circle
        ax2.axis('equal')

        # Insert dynamic data into the center without borders
        ax2.text(0, 0, 'Payment\n Progress:\n', ha='center', va='center', fontsize=30, fontname='Tajawal')
        ax2.text(0, -0.2, str(int(round(payment_progress * 100))) + '%', ha='center', va='center', fontsize=40, fontname='Tajawal')

        # Save the plot to a BytesIO object
        buf2 = io.BytesIO()
        plt.savefig(buf2, format='png', bbox_inches='tight')
        buf2.seek(0)

        # Define the position and size of the image for the second chart
        left2 = 10129900
        top2 = 2970939
        width2 = 2639875
        height2 = 2551024

        # Add the image to the slide from BytesIO object
        new_slide.shapes.add_picture(buf2, left2, top2, width2, height2)
        #-------------------------------------------------------------------------------------------------------------------------------------

        # Define the shape properties to detect
        shape_type = 17  # TEXT_BOX
        position = (1777610, 8048742)
        height = 407543
        width = 3241031

        # Today's date in the desired format
        today_date = datetime.now().strftime("%d %B %Y")

        # Check only the first slide
        slide = ppt.slides[0]
        for shape in slide.shapes:
            if (shape.shape_type == shape_type and
                shape.left == position[0] and
                shape.top == position[1] and
                shape.height == height and
                shape.width == width):

                # Fill the shape's text with today's date
                text_frame = shape.text_frame
                text_frame.clear()  # Clear existing text

                # Directly set the text of the textbox
                text_frame.text = today_date

                # Set font properties for the entire text frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Tajawal'
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(255, 255, 255)  # White color

                # Set vertical alignment to middle
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


          # Save the updated PowerPoint file
    updated_pptx_path = "updated_presentation.pptx"
    ppt.save(updated_pptx_path)
    return updated_pptx_path

st.title("ADF team Project Cards - Excel to PowerPoint Automation")


# Button to show/hide instructions

st.header("Instructions to use this App effectively:")
st.markdown("""
1. Make sure the '07. ADF Automation' folder from SharePoint has been added as a shortcut on your PC via OneDrive (select 'Add Shortcut to OneDrive' from the SharePoint site). Then, select those files while uploading.
2. Ideally, avoid as much blank data as possible, and avoid hiding rows, unless absolutely necessary.
3. Make sure at least all projects from CM tab are not hidden in the PM tab and vice versa.
4. The project PPT images need to be zipped into a folder, with file names matching corresponding project names from CM tab. The ZIP file does not need to be modified after, even if you need to exclude a project. Just mark the 'Include in PPT' column in CM tab as 'no'.
5. If you face any issues or have questions, reach out to: Danish Memon
""")

# File uploaders
excel_file = st.file_uploader("Select Project Cards Excel File", type=["xlsx"])
pptx_file = st.file_uploader("Select PowerPoint Cover Template", type=["pptx"])
image_zip_file = st.file_uploader("Select Images Zip File", type=["zip"])

if excel_file and pptx_file and image_zip_file:
    # Read Excel file directly from the uploaded file
    excel_data = io.BytesIO(excel_file.getbuffer())

    # Read PowerPoint file directly from the uploaded file
    pptx_data = io.BytesIO(pptx_file.getbuffer())

    # Extract images from the uploaded zip file
    image_folder_path = "extracted_images"
    if not os.path.exists(image_folder_path):
        os.makedirs(image_folder_path)

    image_files = extract_images_from_zip(io.BytesIO(image_zip_file.getbuffer()), image_folder_path)

    # Count the number of images in the folder
    num_images = count_images_in_folder(image_folder_path)

    # Process files and overwrite the uploaded PPT file
    updated_pptx_path = read_excel_and_write_to_pptx(excel_data, pptx_data, image_folder_path)

    # Load the updated presentation
    with open(updated_pptx_path, "rb") as f:
        output_pptx = io.BytesIO(f.read())

    st.success(f"File updated successfully! Download below")
    st.download_button(
        label="Download Updated PowerPoint",
        data=output_pptx,
        file_name="ADF Season 2025 Project Cards - Updated.pptx"
    )
