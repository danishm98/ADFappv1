#import tkinter as tk
#from tkinter import filedialog, messagebox
#import pandas as pd
import streamlit as st
#import python-pptx
from pptx import Presentation
import os
#import python-pptx as pptx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
#import pptx
#import pandas as pd
from pandas import read_excel
from openpyxl import load_workbook
from pptx.util import Inches
#import openpyxl
#import matplotlib
from matplotlib.pyplot import savefig, subplots
import io
from PIL import Image
from pptx.enum.chart import XL_LABEL_POSITION
#import pandas as pd
from IPython.display import display
#import ipywidgets as widgets
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR

def read_excel_and_write_to_pptx(excel_path, pptx_path):
    df = read_excel(excel_path)
    ppt = Presentation(pptx_path)
    slide_layout = ppt.slide_layouts[0]
    
    
    # Define properties for the first table (table1)
    table1_position = (304800, 1125677)
    table1_height = 1166985
    table1_width = 12173635
    table1_rows = 2
    table1_columns = 10
    first_row_color = RGBColor(21, 42, 93)  # RGB(21, 42, 93)
    second_row_first_column_color = RGBColor(233, 245, 245)  # RGB(233, 245, 245)
    
    # Insert the new slide immediately after slide 3 (position 4)
    slide_layout = ppt.slides[1].slide_layout  # Use Slide 3's layout
    new_slide = ppt.slides.add_slide(slide_layout)
    xml_slides = ppt.slides._sldIdLst  # Access the low-level XML structure of slides
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])  # Remove the slide that was just added
    xml_slides.insert(3, slides[-1])  # Insert the slide at position 4 (after slide 3)
    
    # Add the first table to the new slide
    table1 = new_slide.shapes.add_table(
        rows=table1_rows,
        cols=table1_columns,
        left=table1_position[0],
        top=table1_position[1],
        width=table1_width,
        height=table1_height
    )
    
    # Reduce row heights by 15%
    row_heights = [int(table1_height * 0.45), int(table1_height * 0.45)]  # 15% less than original height
    for i, row in enumerate(table1.table.rows):
        row.height = row_heights[i]
    
    # Set headers for the first row
    headers = [
        "Project Name",
        "Project Status",
        "Design Status",
        "Construction Start Date",
        "Target Completion Date",
        "Forecast Completion Date",
        "Overall Progress",
        "Current Project Cost",
        "Forecast to Complete",
        "Cost / m2"
    ]
    
    second_row_second_columns = table1.table.rows[1]
    for i, cell in enumerate(second_row_second_columns.cells):
        if i == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(233, 245, 245)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Style the first row with the given color and white, bold text
    first_row = table1.table.rows[0]
    for i, cell in enumerate(first_row.cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = first_row_color  # Set the background color for first row cells
        cell.text_frame.text = headers[i]  # Add header text to the cell
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
        cell.text_frame.paragraphs[0].font.bold = True  # Make text bold
        cell.text_frame.paragraphs[0].font.size = Pt(13)  # Optional: Adjust the font size if needed
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center align the text
        cell.text_frame.paragraphs[0].font.name = 'Tajawal'
        cell.text_frame.margin_top = Pt(0)  # Set top margin to 0
        cell.text_frame.margin_bottom = Pt(0)  # Set bottom margin to 0
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center align the text
    
    # Set column widths
    column_widths = [1486150, 1219005, 1172300, 1159825, 1409073, 1414125, 841825, 1369950, 1107637, 1029837]
    for i, width in enumerate(column_widths):
        table1.table.columns[i].width = width
    
    
    # Define properties for the second table (table2)
    table2_position = (304800, 2271153)
    table2_height = 2287085
    table2_width = 4784621
    table2_first_row_color = RGBColor(21, 42, 93)  # RGB(21, 42, 93)
    table2_second_to_sixth_row_color = RGBColor(233, 245, 245)  # RGB(233, 245, 245)
    
    # Add the second table to the new slide
    table2 = new_slide.shapes.add_table(
        rows=6,
        cols=3,
        left=table2_position[0],
        top=table2_position[1],
        width=table2_width,
        height=table2_height
    )
    
    # Set first row merged with the text "Project Key Stats"
    first_row_table2 = table2.table.rows[0]
    for i, cell in enumerate(first_row_table2.cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(21, 42, 93)    # Set the background color for first row cells
    
    # Merge first row cells (if necessary)
    for i in range(2, 3):
        table2.table.cell(0, i).merge(table2.table.cell(0, 0))
    first_row_table2.cells[0].text_frame.text = "Project Key Stats"  # Add header text
    first_row_table2.cells[0].text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)  # White text color
    first_row_table2.cells[0].text_frame.paragraphs[0].font.bold = True  # Make text bold
    first_row_table2.cells[0].text_frame.paragraphs[0].font.size = Pt(13)  # Font size
    first_row_table2.cells[0].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
    first_row_table2.cells[0].text_frame.paragraphs[0].font.name = 'Tajawal'
    
    # Set the background color for the first two columns (rows 2 through 6)
    for row_idx in range(1, 6):
        for col_idx in range(3):  # For first two columns
            if col_idx == 2:
              cell = table2.table.cell(row_idx, col_idx)
              cell.fill.solid()
              cell.fill.fore_color.rgb = RGBColor(255,255,255)
            else:
              cell = table2.table.cell(row_idx, col_idx)
              cell.fill.solid()
              cell.fill.fore_color.rgb = table2_second_to_sixth_row_color
    
    
    # Set merged rows 2-4 with respective text
    merged_texts = ["Site Area", "Built Up Area", "Project Configuration"]
    for i, text in enumerate(merged_texts):
        cell1 = table2.table.cell(i + 1, 0)
        cell2 = table2.table.cell(i + 1, 1)
        cell1.merge(cell2)  # Merge cells in the first two columns
        cell1.text_frame.text = text  # Set the text
        cell1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
        cell1.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to Pt(12)
        cell1.text_frame.paragraphs[0].font.bold = True  # Bold text
        cell1.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
        cell1.text_frame.paragraphs[0].font.name = 'Tajawal'
    
    # Add merged text for "Space Programme"
    cell5 = table2.table.cell(4, 0)
    cell6 = table2.table.cell(5, 0)
    cell5.merge(cell6)  # Merge rows 5 and 6 in column 1
    cell5.text_frame.text = "Space Programme"  # Add text
    cell5.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
    cell5.text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to Pt(12)
    cell5.text_frame.paragraphs[0].font.bold = True  # Bold text
    cell5.text_frame.paragraphs[0].font.name = 'Tajawal'
    
    # Add text for "FOH" and "BOH"
    table2.table.cell(4, 1).text_frame.text = "FOH"  # Row 5, Column 2
    table2.table.cell(5, 1).text_frame.text = "BOH"  # Row 6, Column 2
    table2.table.cell(4, 1).text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to 12
    table2.table.cell(5, 1).text_frame.paragraphs[0].font.size = Pt(10)  # Set font size to 12
    table2.table.cell(4, 1).text_frame.paragraphs[0].font.bold = True
    table2.table.cell(5, 1).text_frame.paragraphs[0].font.bold = True
    table2.table.cell(4,1).text_frame.paragraphs[0].font.name = 'Tajawal'
    table2.table.cell(5,1).text_frame.paragraphs[0].font.name = 'Tajawal'
    
    # Adjust column widths for table2
    table2.table.columns[0].width = 999725
    table2.table.columns[1].width = 999725
    table2.table.columns[2].width = 2785175
    
    # Add the third table (table3) with 10 rows
    table3_position = (5344484, 3087191)
    table3_height = 2324670
    table3_width = 2384824
    table3_rows = 10
    table3_columns = 2
    
    table3 = new_slide.shapes.add_table(
        rows=table3_rows,
        cols=table3_columns,
        left=table3_position[0],
        top=table3_position[1],
        width=table3_width,
        height=table3_height
    )
    
    
    # Set background color and black non-bold text for Table 3
    for row in range(table3_rows):
        for col in range(table3_columns):
            cell = table3.table.cell(row, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set cell background color to white
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
            cell.text_frame.paragraphs[0].font.bold = False  # Non-bold text
            cell.text_frame.paragraphs[0].font.size = Pt(12)  # Set font size to 12
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align the text
    
    # Draw dashed borders around the table
    border_color = RGBColor(59, 134, 134)
    line_width = Pt(0.65)  # Smaller line width for more frequent dashes
    
    # Top border
    line = new_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, table3_position[0], table3_position[1], table3_position[0] + table3_width, table3_position[1]
    )
    line.line.color.rgb = border_color
    line.line.width = line_width
    line.line.dash_style = MSO_LINE.DASH
    
    # Bottom border
    line = new_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, table3_position[0], table3_position[1] + table3_height, table3_position[0] + table3_width, table3_position[1] + table3_height
    )
    line.line.color.rgb = border_color
    line.line.width = line_width
    line.line.dash_style = MSO_LINE.DASH
    
    # Left border
    line = new_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, table3_position[0], table3_position[1], table3_position[0], table3_position[1] + table3_height
    )
    line.line.color.rgb = border_color
    line.line.width = line_width
    line.line.dash_style = MSO_LINE.DASH
    
    # Right border
    line = new_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, table3_position[0] + table3_width, table3_position[1], table3_position[0] + table3_width, table3_position[1] + table3_height
    )
    line.line.color.rgb = border_color
    line.line.width = line_width
    line.line.dash_style = MSO_LINE.DASH
    # Textboxes (Textbox1 and Textbox2)
    # First Textbox
    textbox1 = new_slide.shapes.add_textbox(
        left=7900000,
        top=2271425,
        width=4590000,
        height=348040
    )
    textbox1.fill.solid()
    textbox1.fill.fore_color.rgb = RGBColor(21, 42, 93)  # Set background color
    text_frame1 = textbox1.text_frame
    text_frame1.text = "Current Construction & Payment Progress"
    text_frame1.paragraphs[0].font.size = Pt(13)
    text_frame1.paragraphs[0].font.bold = True
    text_frame1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    text_frame1.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame1.paragraphs[0].font.name = 'Tajawal'
    
    # Second Textbox
    textbox2 = new_slide.shapes.add_textbox(
        left=5337406,
        top=2271425,
        width=2375000,
        height=348040
    )
    textbox2.fill.solid()
    textbox2.fill.fore_color.rgb = RGBColor(21, 42, 93)  # Set background color
    text_frame2 = textbox2.text_frame
    text_frame2.text = "Project Cost & Cost Analysis"
    text_frame2.paragraphs[0].font.size = Pt(13)
    text_frame2.paragraphs[0].font.bold = True
    text_frame2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    text_frame2.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame2.paragraphs[0].font.name = 'Tajawal'
    
    # Shape 11: Table
    shape_11 = new_slide.shapes.add_table(2, 1, Inches(7629525 / 914400), Inches(6108048 / 914400), Inches(4852309 / 914400), Inches(1179842 / 914400)).table
    # Set row heights (30% for row 1 and 70% for row 2)
    total_height = int(Inches(1179842 / 914400))
    shape_11.rows[0].height = int(total_height * 0.26)
    shape_11.rows[1].height = int(total_height * 0.74)
    shape_11.cell(0, 0).text = "Remaining Items"
    # Set the first row background color to RGBColor(21, 42, 93) with white bold text
    shape_11.cell(0, 0).fill.solid()
    shape_11.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
    cell_11_text_frame = shape_11.cell(0, 0).text_frame
    cell_11_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    cell_11_text_frame.paragraphs[0].font.bold = True
    cell_11_text_frame.paragraphs[0].font.size = Pt(13)
    cell_11_text_frame.paragraphs[0].font.name = "Tajawal"
    # Set the second row background color to RGBColor(233, 245, 245)
    shape_11.cell(1, 0).fill.solid()
    shape_11.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
    shape_11.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
    
    # Shape 12: Table
    shape_12 = new_slide.shapes.add_table(3, 2, Inches(5342195 / 914400), Inches(6111314 / 914400), Inches(2134930 / 914400), Inches(1173428 / 914400)).table
    shape_12.cell(0, 0).text = "Estimated Total Required Avg. Daily Resources"
    # Merge the first row
    cell_12_0_0 = shape_12.cell(0, 0)
    cell_12_0_1 = shape_12.cell(0, 1)
    cell_12_0_0.merge(cell_12_0_1)
    # Set the first row background color to RGBColor(21, 42, 93) with white bold text
    cell_12_0_0.fill.solid()
    cell_12_0_0.fill.fore_color.rgb = RGBColor(21, 42, 93)
    cell_12_text_frame = cell_12_0_0.text_frame
    cell_12_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    cell_12_text_frame.paragraphs[0].font.bold = True
    cell_12_text_frame.paragraphs[0].font.size = Pt(13)
    cell_12_text_frame.paragraphs[0].font.name = "Tajawal"
    cell_12_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Set the second row background color to RGBColor(233, 245, 245)
    shape_12.cell(1, 0).fill.solid()
    shape_12.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
    shape_12.cell(1, 1).fill.solid()
    shape_12.cell(1, 1).fill.fore_color.rgb = RGBColor(233, 245, 245)
    # Insert the values "Daily Manpower" and "Daily Machinery" in the second row with black bold point 10.5 text
    shape_12.cell(1, 0).text = "Daily Manpower"
    shape_12.cell(1, 1).text = "Daily Machinery"
    cell_12_text_frame_1 = shape_12.cell(1, 0).text_frame
    cell_12_text_frame_1.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text color
    cell_12_text_frame_1.paragraphs[0].font.bold = True
    cell_12_text_frame_1.paragraphs[0].font.size = Pt(12)
    cell_12_text_frame_2 = shape_12.cell(1, 1).text_frame
    cell_12_text_frame_2.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text color
    cell_12_text_frame_2.paragraphs[0].font.bold = True
    cell_12_text_frame_2.paragraphs[0].font.size = Pt(12)
    cell_12_text_frame_1.paragraphs[0].font.name = "Tajawal"
    cell_12_text_frame_2.paragraphs[0].font.name = "Tajawal"
    cell_12_text_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
    cell_12_text_frame_2.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
    # Set the last row background color to white
    shape_12.cell(2, 0).fill.solid()
    shape_12.cell(2, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
    shape_12.cell(2, 1).fill.solid()
    shape_12.cell(2, 1).fill.fore_color.rgb = RGBColor(233, 245, 245)
    # Adjust row heights (30% for row 1 and 70% for row 2)
    total_height = Inches(1173428 / 914400)
    cell_12_text_frame_1.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
    
    # Shape 13: Table
    shape_13 = new_slide.shapes.add_table(1, 2, Inches(5333250 / 914400), Inches(2673352 / 914400), Inches(2362950 / 914400), Inches(365760 / 914400)).table
    shape_13.cell(0, 0).text = "Current Project Cost"
    shape_13.cell(0, 0).fill.solid()
    shape_13.cell(0, 0).fill.fore_color.rgb = RGBColor(59, 134, 134)
    # Make cell (0,0)'s text bold and set font size to 12
    cell_00_text_frame = shape_13.cell(0, 0).text_frame
    cell_00_text_frame.paragraphs[0].font.bold = True
    cell_00_text_frame.paragraphs[0].font.size = Pt(11)
    cell_00_text_frame.paragraphs[0].font.name = "Tajawal"
    # Set the second column's color
    shape_13.cell(0, 1).fill.solid()
    shape_13.cell(0, 1).fill.fore_color.rgb = RGBColor(211,236,236)
    # Set the column widths to 65% and 35% for the first two columns respectively
    total_width = int(Inches(2362950 / 914400))
    first_col_width = int(total_width * 0.65)
    second_col_width = int(total_width * 0.35)
    shape_13.columns[0].width = first_col_width
    shape_13.columns[1].width = second_col_width
    shape_13.rows[0].height =int(shape_13.rows[0].height*0.9)
    
    # Shape 15: Table
    shape_15 = new_slide.shapes.add_table(2, 1, Inches(5332081 / 914400), Inches(5437224 / 914400) , Inches(7169541 / 914400), Inches(634572 / 914400)).table
    shape_15.cell(0, 0).text = "Basis of Project Cost"
    # Set the first row background color to RGBColor(21, 42, 93) with white bold text
    shape_15.cell(0, 0).fill.solid()
    shape_15.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
    cell_15_text_frame = shape_15.cell(0, 0).text_frame
    cell_15_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    cell_15_text_frame.paragraphs[0].font.bold = True
    cell_15_text_frame.paragraphs[0].font.size = Pt(13)
    cell_15_text_frame.paragraphs[0].font.name = "Tajawal"
    # Set the second row background color to white
    shape_15.cell(1, 0).fill.solid()
    shape_15.cell(1, 0).fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # table 20---------------------
    values = [
        "Design & Supervision",
        "Facilitating Works",
        "Substructure",
        "Superstructure",
        #"Fittings, Furnishings & Equipment",
        "Specialties",
        "Internal Finishes",
        "Services Features",
       # "External Works",
       # "L. Hardscape",
       # "L. Softscape",
       # "Inf. Utility Building",
       # "Inf. Networks",
       # "Main Contractor Preliminaries",
       # "Main Contractor Overhead & Profit",
       # "Client Direct Procurement",
       # "Design Contingencies",
       # "Construction Contingencies",
        "Forecast Construction Spend 2025:"
    ]
    # Add the table with specified position and size
    num_rows = len(values)
    num_cols = 2
    shape_20 = new_slide.shapes.add_table(num_rows, num_cols, Inches(5344484 / 914400), Inches(3087191 / 914400) + Inches(0.03), Inches(2384824 / 914400), Inches(2324670 / 914400)).table
    
    # Set the first column width to 70% of the total width
    total_width = int(Inches(2384824 / 914400))
    first_col_width = int(total_width * 0.70)
    second_col_width = int(total_width * 0.30)
    shape_20.columns[0].width = first_col_width
    shape_20.columns[1].width = second_col_width
    
    # Set all columns color to RGBColor(233, 245, 245)
    for row in shape_20.rows:
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(233, 245, 245)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Set the values in the first column and set text size to Pt(10.5) and not bold
    for i, value in enumerate(values):
        shape_20.cell(i, 0).text = value
        cell_text_frame = shape_20.cell(i, 0).text_frame
        cell_text_frame.paragraphs[0].font.size = Pt(11)
        cell_text_frame.paragraphs[0].font.bold = False
        cell_text_frame.paragraphs[0].font.name = "Tajawal"
    
    # Ensure the first row is not styled as a heading/header and has normal text like the other cells
    for i in range(num_cols):
        cell_text_frame = shape_20.cell(0, i).text_frame
        cell_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set text color to black
        cell_text_frame.paragraphs[0].font.bold = False
        cell_text_frame.paragraphs[0].font.size = Pt(11)
        cell_text_frame.paragraphs[0].font.name = "Tajawal"
    
    # Merge the cells in the last row
    shape_20.cell(num_rows - 1, 0).merge(shape_20.cell(num_rows - 1, 1))
    
    # Set the background color to RGB(59, 134, 134) and text to white bold for the last row
    last_row_cell = shape_20.cell(num_rows - 1, 0)
    last_row_cell.fill.solid()
    last_row_cell.fill.fore_color.rgb = RGBColor(29, 88, 137)
    last_row_text_frame = last_row_cell.text_frame
    last_row_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    last_row_text_frame.paragraphs[0].font.bold = True
    last_row_text_frame.paragraphs[0].font.name = "Tajawal"
    
    # Adjust row heights to increase the last row's height by 30% and compensate in other rows accordingly
    total_height = Inches(2324670 / 914400)
    current_row_height = total_height / num_rows
    new_last_row_height = current_row_height * 1.3
    height_difference = new_last_row_height - current_row_height
    new_other_row_height = (total_height - new_last_row_height) / (num_rows - 1)
    
    for i in range(num_rows - 1):
        shape_20.rows[i].height = int(new_other_row_height)
    
    shape_20.rows[num_rows - 1].height = int(new_last_row_height)
    # Shape 21: Table
    shape_21 = new_slide.shapes.add_table(2, 1, Inches(5342195 / 914400), Inches(7390134 / 914400), Inches(7136239 / 914400), Inches(1015802.4 / 914400)).table
    shape_21.cell(0, 0).text = "Risk Assessment"
    # Set the first row background color to RGBColor(21, 42, 93) with white bold text
    shape_21.cell(0, 0).fill.solid()
    shape_21.cell(0, 0).fill.fore_color.rgb = RGBColor(21, 42, 93)
    total_height=Inches(1015802.4 / 914400)
    cell_21_text_frame = shape_21.cell(0, 0).text_frame
    cell_21_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text color
    cell_21_text_frame.paragraphs[0].font.bold = True
    cell_21_text_frame.paragraphs[0].font.size = Pt(13)
    cell_21_text_frame.paragraphs[0].font.name = "Tajawal"
    cell_21_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Left align the text
    # Set the second row background color to RGBColor(233, 245, 245)
    shape_21.cell(1, 0).fill.solid()
    shape_21.cell(1, 0).fill.fore_color.rgb = RGBColor(233, 245, 245)
    shape_21.rows[0].height = int(total_height * 0.3)
    shape_21.rows[1].height = int(total_height * 0.7)
    shape_21.rows[1].height = int(shape_21.rows[1].height * 1.3)
    
    
    # Chart data
    # Data for the donut chart
    sizes = [93, 7]
    labels = ['Completed', 'In Progress']
    colors = ['#0aa57f', '#1d5889']  # RGB: (10, 165, 127) and (29, 88, 137)
    
    # Create a figure and axis with equal aspect ratio to avoid squeezing
    fig, ax = subplots(figsize=(8, 8))
    
    # Create the donut chart with thicker white borders for slices
    wedges, texts, autotexts = ax.pie(sizes, colors=colors, autopct=lambda p: '{:.1f}%'.format(p) if p > 7 else '',
                                      startangle=90, wedgeprops=dict(width=0.3, edgecolor='white', linewidth=3), pctdistance=0.85, textprops=dict(color='white', fontsize=28))
    
    # Add legends at the bottom of the chart in one line without borders or shadows and make them bigger
    ax.legend(wedges, labels, loc="upper center", bbox_to_anchor=(0.5, -0.1), frameon=False, ncol=2, fontsize=25)
    
    # Equal aspect ratio ensures that pie is drawn as a circle
    ax.axis('equal')
    
    # Insert dynamic data into the center without borders
    ax.text(0, 0, 'Construction\n Progress:\n7%', ha='center', va='center', fontsize=25)
    
    # Save the plot to a BytesIO object
    buf = io.BytesIO()
    savefig(buf, format='png', bbox_inches='tight')
    buf.seek(0)
    
    #chart_shape.left = 7470000
    #chart_shape.top = 2650000
    
    # Define the position and size of the image
    left = 7884225
    top = 2964434
    width = 2245675
    height = 2355166  # Adjusted height to maintain aspect ratio
    
    # Add the image to the slide from BytesIO object
    new_slide.shapes.add_picture(buf, left, top, width, height)
    
    
    
    left = 10284825
    top = 2964434
    height: 2355166
    width: 2245675
    new_slide.shapes.add_picture(buf, left, top, width, height)
    

    ppt.save(pptx_path)

st.title("ADF team Project Cards - Excel to PowerPoint Converter")

# File uploaders
excel_file = st.file_uploader("Select Excel File", type=["xlsx"])
pptx_file = st.file_uploader("Select PowerPoint File", type=["pptx"])

if excel_file and pptx_file:
    # Save uploaded files to a temporary directory
    with open("temp_excel.xlsx", "wb") as f:
        f.write(excel_file.getbuffer())
    with open("temp_pptx.pptx", "wb") as f:
        f.write(pptx_file.getbuffer())

    # Process files and overwrite the uploaded PPT file
    read_excel_and_write_to_pptx("temp_excel.xlsx", "temp_pptx.pptx")
    
    st.success("File updated successfully!")
    st.download_button(
        label="Download Updated PowerPoint",
        data=open("temp_pptx.pptx", "rb").read(),
        file_name="Updated_Presentation.pptx"
    )
